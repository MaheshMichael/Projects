from pptx import Presentation
import io

def create_pptx(summary_dict, template_path='pptx_templates/EYtemplate-simple.pptx'):
    """
    Creates a PowerPoint presentation from a summary dictionary and a template.
    Returns a BytesIO stream for downloading.
    """
    prs = Presentation(template_path)
    
    # Updating slides with summary_dict data
    slides_data = summary_dict.get("slides", [])
    
    slide = prs.slides[0]
    slide.placeholders[0].text = slides_data[0]["title"]
    slide.placeholders[1].text = slides_data[0]["subtitle"]

    slide = prs.slides[1]
    slide.placeholders[0].text = slides_data[1]["title"]


    slide = prs.slides[2]
    slide.placeholders[0].text = slides_data[2]["title"]
    slide.placeholders[16].text = slides_data[2]["content"]

    slide = prs.slides[3]
    slide.placeholders[0].text = slides_data[3]["title"]

    slide = prs.slides[4]
    slide.placeholders[0].text = slides_data[4]["title"]
    slide.placeholders[16].text = slides_data[4]["content"]

    slide = prs.slides[5]
    slide.placeholders[0].text = slides_data[5]["title"]
    slide.placeholders[16].text = "\n".join(slides_data[5]["content"]["placeholder_1"])
    slide.placeholders[15].text = "\n".join(slides_data[5]["content"]["placeholder_2"])

    slide = prs.slides[6]
    slide.placeholders[4294967295].text = "\n".join(slides_data[6]["content"])

    slide = prs.slides[7]
    slide.placeholders[0].text = slides_data[7]["title"]
    slide.placeholders[15].text = "\n".join(slides_data[7]["content"])
    
    # Save to BytesIO instead of file system
    pptx_stream = io.BytesIO()
    prs.save(pptx_stream)  # Save the presentation to the BytesIO stream
    pptx_stream.seek(0)  # Reset stream position for reading
    
    return pptx_stream
